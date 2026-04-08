"""
Guia de Implantação do Sistema de Pagamento — PsicoDoc
Tutorial passo a passo para leigos configurarem Stripe + Supabase
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paleta de cores ──────────────────────────────────────────────
AZUL       = RGBColor(0x4F, 0x7C, 0xAC)
VERDE      = RGBColor(0x2E, 0x9E, 0x6B)
AMBAR      = RGBColor(0xE8, 0xA8, 0x38)
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_F    = RGBColor(0xF4, 0xF6, 0xFA)
TEXTO      = RGBColor(0x1E, 0x2A, 0x38)
CINZA_M    = RGBColor(0x64, 0x74, 0x8B)
AZUL_ESC   = RGBColor(0x1E, 0x40, 0x6E)
ROXO       = RGBColor(0x63, 0x4D, 0xBF)
LARANJA    = RGBColor(0xF5, 0x5C, 0x1B)
VERDE_ESC  = RGBColor(0x14, 0x6B, 0x3A)

W = Inches(13.333)
H = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────

def nova_apresentacao():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def slide_branco(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, left, top, width, height, fill_rgb):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def txt(slide, text, left, top, width, height,
        size=18, bold=False, color=TEXTO, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def txt_multi(slide, lines, left, top, width, height,
              size=16, bold=False, color=TEXTO, align=PP_ALIGN.LEFT,
              bold_first=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for i, line in enumerate(lines):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold or (bold_first and i == 0)
        run.font.color.rgb = color
    return txBox

def card(slide, left, top, width, height, bg, titulo, corpo_lines,
         titulo_size=18, corpo_size=15, titulo_cor=BRANCO, corpo_cor=TEXTO,
         corpo_bg=None):
    rect(slide, left, top, width, height, bg)
    if corpo_bg:
        rect(slide, left, top + 0.55, width, height - 0.55, corpo_bg)
    txt(slide, titulo, left + 0.15, top + 0.1, width - 0.3, 0.45,
        size=titulo_size, bold=True, color=titulo_cor)
    txt_multi(slide, corpo_lines, left + 0.15, top + 0.62,
              width - 0.3, height - 0.75, size=corpo_size, color=corpo_cor)

def numero_slide(slide, n, total):
    txt(slide, f"{n} / {total}", 12.5, 7.1, 0.7, 0.3,
        size=11, color=CINZA_M, align=PP_ALIGN.RIGHT)

def cabecalho_faixa(slide, titulo, subtitulo="", cor=AZUL):
    rect(slide, 0, 0, 13.333, 1.4, cor)
    txt(slide, titulo, 0.4, 0.15, 12, 0.7,
        size=28, bold=True, color=BRANCO)
    if subtitulo:
        txt(slide, subtitulo, 0.4, 0.82, 12, 0.5,
            size=16, color=BRANCO, italic=True)

def rodape(slide, msg, cor=CINZA_M):
    rect(slide, 0, 7.1, 13.333, 0.4, CINZA_F)
    txt(slide, msg, 0.3, 7.12, 12, 0.3, size=11, color=cor)

# ════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════

def slide_capa(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, AZUL_ESC)

    # Faixa decorativa lateral
    rect(sl, 0, 0, 0.25, 7.5, VERDE)

    # Ícone visual (retângulo simulando cartão de crédito)
    rect(sl, 9.8, 2.2, 3.0, 1.9, ROXO)
    rect(sl, 9.8, 2.2, 3.0, 0.55, RGBColor(0x4A, 0x35, 0x9E))
    rect(sl, 10.0, 3.2, 0.7, 0.45, AMBAR)
    rect(sl, 10.85, 3.2, 1.6, 0.45, RGBColor(0x9A, 0x8F, 0xD9))

    txt(sl, "💳", 10.55, 2.3, 1.2, 0.7, size=28, color=BRANCO, align=PP_ALIGN.CENTER)

    txt(sl, "PsicoDoc", 0.6, 1.0, 9, 0.7,
        size=20, bold=False, color=VERDE, align=PP_ALIGN.LEFT)
    txt(sl, "Guia de Implantação do", 0.6, 1.7, 10, 0.9,
        size=36, bold=False, color=BRANCO, align=PP_ALIGN.LEFT)
    txt(sl, "Sistema de Pagamento", 0.6, 2.55, 10, 0.9,
        size=42, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)

    txt(sl, "Como configurar cobranças recorrentes no cartão\ne aceitar PIX no seu app — passo a passo para leigos",
        0.6, 3.65, 9.5, 1.2, size=18, color=RGBColor(0xB0, 0xC8, 0xE8),
        align=PP_ALIGN.LEFT, wrap=True)

    # Chips de tecnologia
    for i, (label, cor) in enumerate([
        ("Stripe", ROXO),
        ("Supabase", VERDE),
        ("Vercel", AZUL),
    ]):
        x = 0.6 + i * 2.3
        rect(sl, x, 5.2, 2.0, 0.45, cor)
        txt(sl, label, x, 5.2, 2.0, 0.45,
            size=15, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

    txt(sl, "Nenhum conhecimento técnico necessário",
        0.6, 6.5, 9, 0.4, size=13, color=CINZA_M, italic=True)


def slide_visao_geral(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, CINZA_F)
    cabecalho_faixa(sl, "O que vamos configurar?",
                    "Visão geral de todo o processo", AZUL_ESC)

    blocos = [
        (ROXO,      "1. Criar conta no Stripe",
         ["O Stripe é a empresa que processa\nos pagamentos com cartão",
          "Gratuito para começar — você só\npaga quando receber"]),
        (AZUL,      "2. Criar o produto e preço",
         ["Cadastrar o plano R$ 97/mês\ncomo produto recorrente no Stripe",
          "Gerar o Price ID que o sistema usa\npara cobrar automaticamente"]),
        (VERDE,     "3. Configurar o Supabase",
         ["Rodar a migração SQL que adiciona\nas colunas de assinatura",
          "Configurar URL do site e\nconfirmação de e-mail"]),
        (AMBAR,     "4. Variáveis no Vercel",
         ["Adicionar as chaves secretas\ndo Stripe e Supabase no Vercel",
          "Fazer redeploy para ativar\nas novas configurações"]),
    ]

    for i, (cor, titulo, linhas) in enumerate(blocos):
        col = i % 2
        row = i // 2
        x = 0.4 + col * 6.4
        y = 1.65 + row * 2.7
        card(sl, x, y, 6.1, 2.5, cor, titulo, linhas,
             corpo_bg=BRANCO, corpo_cor=TEXTO)

    rodape(sl, "Tempo estimado: 30 a 40 minutos   |   Você não precisa saber programar")
    numero_slide(sl, 2, 13)


def slide_stripe_conta(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, CINZA_F)
    cabecalho_faixa(sl, "Passo 1 — Criar conta no Stripe",
                    "stripe.com  •  Gratuito para começar", ROXO)

    # Coluna esquerda: passos
    rect(sl, 0.4, 1.6, 7.2, 5.3, BRANCO)
    txt(sl, "O que é o Stripe?", 0.6, 1.7, 7.0, 0.5,
        size=17, bold=True, color=ROXO)
    txt_multi(sl, [
        "O Stripe é o sistema que vai cobrar seus clientes com cartão de crédito.",
        "Ele é usado por empresas como Amazon, Shopify e Nubank.",
        "Você recebe diretamente na sua conta bancária.",
    ], 0.6, 2.2, 7.0, 1.0, size=14, color=TEXTO)

    txt(sl, "Como criar sua conta:", 0.6, 3.3, 7.0, 0.5,
        size=17, bold=True, color=AZUL_ESC)

    passos = [
        ("1", "Acesse stripe.com e clique em Comece agora"),
        ("2", "Preencha seu e-mail, nome e senha"),
        ("3", "Confirme o e-mail que o Stripe vai enviar"),
        ("4", "Preencha os dados do seu negócio (CPF/CNPJ)"),
        ("5", "Adicione sua conta bancária para receber pagamentos"),
    ]
    for i, (num, passo) in enumerate(passos):
        y = 3.85 + i * 0.47
        rect(sl, 0.6, y, 0.38, 0.35, ROXO)
        txt(sl, num, 0.6, y, 0.38, 0.35,
            size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(sl, passo, 1.1, y, 6.3, 0.38, size=13, color=TEXTO)

    # Coluna direita: aviso
    rect(sl, 7.9, 1.6, 5.0, 2.5, RGBColor(0xF0, 0xEB, 0xFC))
    txt(sl, "⚠️  Modo Teste vs Produção", 8.05, 1.7, 4.8, 0.5,
        size=15, bold=True, color=ROXO)
    txt_multi(sl, [
        "O Stripe tem dois modos:",
        "",
        "🧪 Teste — para experimentar sem\n    cobrar ninguém de verdade",
        "",
        "✅ Produção — cobranças reais",
        "",
        "Comece no modo TESTE.\nAtive a produção só quando\nestiver tudo funcionando.",
    ], 8.05, 2.25, 4.75, 2.7, size=13, color=TEXTO)

    rect(sl, 7.9, 4.3, 5.0, 2.6, RGBColor(0xE8, 0xF8, 0xEE))
    txt(sl, "💰  Quanto custa o Stripe?", 8.05, 4.4, 4.8, 0.5,
        size=15, bold=True, color=VERDE_ESC)
    txt_multi(sl, [
        "Sem mensalidade!",
        "",
        "Você paga apenas quando recebe:",
        "• Cartão nacional: 3,49% + R$ 0,39",
        "• Cartão internacional: 5,99%",
        "",
        "O Stripe desconta automático.",
    ], 8.05, 4.92, 4.75, 1.85, size=13, color=TEXTO)

    rodape(sl, "stripe.com  →  clique em Comece agora ou Start now")
    numero_slide(sl, 3, 13)


def slide_stripe_produto(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, CINZA_F)
    cabecalho_faixa(sl, "Passo 2 — Criar o produto e o preço",
                    "Stripe Dashboard  →  Produtos  →  Adicionar produto", ROXO)

    # Caixa de aviso
    rect(sl, 0.4, 1.6, 12.5, 0.65, RGBColor(0xFF, 0xF3, 0xCD))
    txt(sl, "⚠️  Importante: faça isso no modo PRODUÇÃO (não no modo Teste) quando for cobrar de verdade.",
        0.55, 1.68, 12.2, 0.45, size=14, bold=True, color=RGBColor(0x7B, 0x5B, 0x00))

    # Passos lado a lado
    passos = [
        (ROXO, "No painel do Stripe:", [
            "1. Clique em Produtos no menu lateral",
            "2. Clique em + Adicionar produto",
            "3. Preencha o nome: PsicoDoc Mensal",
            "4. Descrição (opcional): Acesso mensal ao sistema",
        ]),
        (AZUL, "Configurar o preço:", [
            "5. Em Preco, selecione Recorrente",
            "6. Valor: R$ 97,00",
            "7. Moeda: BRL (Real Brasileiro)",
            "8. Período: Mensal",
        ]),
        (VERDE, "Salvar e copiar o ID:", [
            "9. Clique em Salvar produto",
            "10. Na página do produto, clique no preço",
            "11. Copie o ID do preco (começa com price_...)",
            "12. Guarde esse código — precisaremos dele!",
        ]),
    ]

    for i, (cor, titulo, linhas) in enumerate(passos):
        x = 0.4 + i * 4.3
        card(sl, x, 2.45, 4.0, 4.35, cor, titulo, linhas,
             corpo_bg=BRANCO, corpo_cor=TEXTO, corpo_size=13)

    # Destaque do Price ID
    rect(sl, 0.4, 6.75, 12.5, 0.55, RGBColor(0xE8, 0xF8, 0xEE))
    txt(sl, "✅  O Price ID tem este formato:  price_1RaBcDeFgHiJkLmNoPqRsTuV   — guarde ele num bloco de notas!",
        0.55, 6.82, 12.2, 0.38, size=13, bold=True, color=VERDE_ESC)

    numero_slide(sl, 4, 13)


def slide_stripe_chaves(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, CINZA_F)
    cabecalho_faixa(sl, "Passo 3 — Copiar as chaves secretas do Stripe",
                    "Stripe Dashboard  →  Desenvolvedores  →  Chaves de API", ROXO)

    # Aviso
    rect(sl, 0.4, 1.6, 12.5, 0.6, RGBColor(0xFF, 0xEB, 0xEB))
    txt(sl, "🔒  As chaves secretas são como senhas do seu sistema de pagamento. NUNCA compartilhe com ninguém.",
        0.55, 1.68, 12.2, 0.4, size=14, bold=True, color=RGBColor(0x9B, 0x1C, 0x1C))

    # Chaves
    chaves = [
        (ROXO, "Chave Publicável\n(Publishable Key)",
         "pk_live_...",
         ["Começa com pk_live_ (produção)\nou pk_test_ (teste)",
          "Não é tão secreta, mas guarde assim mesmo",
          "No nosso app NÃO usamos diretamente"],
         False),
        (AZUL_ESC, "Chave Secreta\n(Secret Key)",
         "sk_live_...",
         ["⭐ Esta é a mais importante!",
          "Começa com sk_live_ (produção)",
          "Coloque no Vercel como STRIPE_SECRET_KEY",
          "NUNCA cole em código ou mande por WhatsApp"],
         True),
    ]

    for i, (cor, titulo, exemplo, linhas, destaque) in enumerate(chaves):
        x = 0.4 + i * 6.3
        rect(sl, x, 2.35, 5.9, 4.2, cor)
        rect(sl, x, 2.35, 5.9, 0.7, RGBColor(
            max(cor[0] - 30, 0), max(cor[1] - 30, 0), max(cor[2] - 30, 0)
        ))
        txt(sl, titulo, x + 0.15, 2.38, 5.6, 0.65,
            size=16, bold=True, color=BRANCO)
        rect(sl, x + 0.15, 3.15, 5.6, 0.5, RGBColor(0x1A, 0x1A, 0x2E))
        txt(sl, exemplo, x + 0.25, 3.2, 5.4, 0.4,
            size=14, bold=True, color=RGBColor(0x7D, 0xFF, 0xD8))
        txt_multi(sl, linhas, x + 0.15, 3.78, 5.6, 2.6,
                  size=13, color=BRANCO)

    # Como acessar
    rect(sl, 0.4, 6.6, 12.5, 0.7, BRANCO)
    txt_multi(sl, [
        "Como encontrar:  Stripe → menu lateral → Desenvolvedores -> Chaves de API → clique no ícone de olho → copie a chave secreta"
    ], 0.55, 6.65, 12.2, 0.5, size=13, color=AZUL_ESC, bold=True)

    numero_slide(sl, 5, 13)


def slide_stripe_webhook(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, CINZA_F)
    cabecalho_faixa(sl, "Passo 4 — Configurar o Webhook do Stripe",
                    "Stripe Dashboard  →  Desenvolvedores  →  Webhooks", ROXO)

    # Explicação do webhook
    rect(sl, 0.4, 1.6, 5.5, 5.6, BRANCO)
    txt(sl, "O que é um webhook?", 0.55, 1.7, 5.2, 0.45,
        size=17, bold=True, color=ROXO)
    txt_multi(sl, [
        "É uma campainha que o Stripe toca",
        "no seu sistema quando algo acontece:",
        "",
        "🔔  Cliente assinou → sistema ativa",
        "🔔  Pagamento recebido → atualiza status",
        "🔔  Pagamento falhou → bloqueia acesso",
        "🔔  Assinatura cancelada → suspende",
        "",
        "Sem o webhook, seu sistema não saberia",
        "se o cliente pagou ou cancelou!",
    ], 0.55, 2.2, 5.2, 4.8, size=14, color=TEXTO)

    # Passos
    rect(sl, 6.2, 1.6, 6.8, 5.6, RGBColor(0xF0, 0xEB, 0xFC))
    txt(sl, "Como configurar:", 6.35, 1.7, 6.5, 0.45,
        size=17, bold=True, color=ROXO)

    passos_wh = [
        ("1", "No Stripe, vá em Desenvolvedores → Webhooks"),
        ("2", "Clique em + Adicionar endpoint"),
        ("3", "URL do endpoint — cole exatamente:"),
        ("", "https://psicogest-one.vercel.app/api/stripe/webhook"),
        ("4", "Em Eventos, clique em + Selecionar eventos"),
        ("5", "Selecione estes 5 eventos:"),
        ("", "• customer.subscription.created"),
        ("", "• customer.subscription.updated"),
        ("", "• customer.subscription.deleted"),
        ("", "• invoice.payment_succeeded"),
        ("", "• invoice.payment_failed"),
        ("6", "Clique em Adicionar endpoint"),
        ("7", "Copie o Signing secret (whsec_...)"),
    ]

    y_base = 2.22
    for num, texto in passos_wh:
        if num:
            rect(sl, 6.35, y_base, 0.32, 0.28, ROXO)
            txt(sl, num, 6.35, y_base, 0.32, 0.28,
                size=11, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
            txt(sl, texto, 6.75, y_base, 6.1, 0.28, size=12, color=TEXTO)
        else:
            rect(sl, 6.55, y_base, 6.3, 0.28, RGBColor(0x1A, 0x1A, 0x2E))
            txt(sl, texto, 6.6, y_base, 6.2, 0.28,
                size=11, bold=True, color=RGBColor(0x7D, 0xFF, 0xD8))
        y_base += 0.35

    rodape(sl, "Guarde o Signing Secret (whsec_...) — vai para o Vercel como STRIPE_WEBHOOK_SECRET")
    numero_slide(sl, 6, 13)


def slide_supabase_migracao(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, CINZA_F)
    cabecalho_faixa(sl, "Passo 5 — Atualizar o banco de dados (Supabase)",
                    "Supabase Dashboard  →  SQL Editor", VERDE)

    txt(sl, "Vamos adicionar as colunas de assinatura ao banco de dados rodando um script SQL.",
        0.4, 1.65, 12.5, 0.45, size=15, color=TEXTO)

    # Instruções
    rect(sl, 0.4, 2.2, 5.2, 4.9, BRANCO)
    txt(sl, "Como fazer:", 0.55, 2.3, 5.0, 0.4,
        size=17, bold=True, color=VERDE_ESC)
    passos = [
        ("1", "Acesse supabase.com e faça login"),
        ("2", "Abra seu projeto PsicoDoc"),
        ("3", "No menu lateral, clique em\nSQL Editor"),
        ("4", "Clique em New query"),
        ("5", "Copie e cole TODO o código\ndo slide seguinte"),
        ("6", "Clique no botão verde Run"),
        ("7", "Deve aparecer Success"),
    ]
    for i, (num, texto) in enumerate(passos):
        y = 2.78 + i * 0.58
        rect(sl, 0.55, y, 0.35, 0.32, VERDE)
        txt(sl, num, 0.55, y, 0.35, 0.32,
            size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(sl, texto, 1.0, y, 4.5, 0.52, size=13, color=TEXTO)

    # O que faz
    rect(sl, 5.9, 2.2, 7.0, 4.9, RGBColor(0xE8, 0xF8, 0xEE))
    txt(sl, "O que esse script faz?", 6.05, 2.3, 6.7, 0.4,
        size=17, bold=True, color=VERDE_ESC)
    txt_multi(sl, [
        "Adiciona novas colunas à tabela de psicólogos:",
        "",
        "📋  plan — se está em trial, ativo ou suspenso",
        "📋  trial_ends_at — data de fim do período grátis",
        "📋  subscription_status — status da assinatura",
        "📋  stripe_customer_id — ID do cliente no Stripe",
        "📋  stripe_subscription_id — ID da assinatura",
        "📋  last_payment_at — data do último pagamento",
        "📋  next_payment_at — data do próximo pagamento",
        "📋  is_admin — indica se é o administrador",
        "",
        "Todos os usuários existentes recebem",
        "automaticamente 14 dias de trial gratuito.",
    ], 6.05, 2.78, 6.7, 4.2, size=13, color=TEXTO)

    numero_slide(sl, 7, 13)


def slide_sql_codigo(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, RGBColor(0x0D, 0x1B, 0x2A))

    txt(sl, "Código SQL — copie e cole no Supabase SQL Editor",
        0.4, 0.15, 12.5, 0.5, size=16, bold=True, color=VERDE)

    codigo = [
        "ALTER TABLE psychologists ADD COLUMN IF NOT EXISTS plan TEXT DEFAULT 'trial'",
        "  CHECK (plan IN ('trial', 'active', 'suspended', 'cancelled'));",
        "",
        "ALTER TABLE psychologists ADD COLUMN IF NOT EXISTS trial_ends_at TIMESTAMPTZ",
        "  DEFAULT (NOW() + INTERVAL '14 days');",
        "",
        "UPDATE psychologists SET trial_ends_at = created_at + INTERVAL '14 days'",
        "WHERE trial_ends_at IS NULL;",
        "",
        "ALTER TABLE psychologists ADD COLUMN IF NOT EXISTS stripe_customer_id TEXT;",
        "ALTER TABLE psychologists ADD COLUMN IF NOT EXISTS stripe_subscription_id TEXT;",
        "ALTER TABLE psychologists ADD COLUMN IF NOT EXISTS subscription_status TEXT DEFAULT 'trialing';",
        "ALTER TABLE psychologists ADD COLUMN IF NOT EXISTS last_payment_at TIMESTAMPTZ;",
        "ALTER TABLE psychologists ADD COLUMN IF NOT EXISTS next_payment_at TIMESTAMPTZ;",
        "ALTER TABLE psychologists ADD COLUMN IF NOT EXISTS is_admin BOOLEAN DEFAULT FALSE;",
        "",
        "-- Marque SEU e-mail como administrador (substitua pelo seu e-mail real):",
        "UPDATE psychologists SET is_admin = true WHERE email = 'SEU-EMAIL@gmail.com';",
    ]

    rect(sl, 0.3, 0.72, 12.7, 6.5, RGBColor(0x0A, 0x14, 0x1E))
    txt_multi(sl, codigo, 0.45, 0.78, 12.4, 6.3,
              size=12, color=RGBColor(0x7D, 0xFF, 0xD8))

    txt(sl, "⚠️  Troque SEU-EMAIL@gmail.com pelo seu e-mail real antes de executar!",
        0.4, 7.1, 12.5, 0.35, size=13, bold=True, color=AMBAR)

    numero_slide(sl, 8, 13)


def slide_supabase_url(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, CINZA_F)
    cabecalho_faixa(sl, "Passo 6 — Corrigir URL e confirmar e-mail no Supabase",
                    "Supabase Dashboard  →  Authentication  →  URL Configuration", VERDE)

    # Box esquerdo: URL
    rect(sl, 0.4, 1.6, 5.9, 5.5, BRANCO)
    txt(sl, "🌐  Configurar URL do site", 0.55, 1.7, 5.7, 0.45,
        size=17, bold=True, color=VERDE_ESC)
    txt(sl, "Problema: os e-mails enviados pelo Supabase apontam para localhost.", 0.55, 2.22, 5.7, 0.55,
        size=13, color=TEXTO, wrap=True)
    txt(sl, "Como corrigir:", 0.55, 2.85, 5.7, 0.4, size=14, bold=True, color=AZUL_ESC)
    passos_url = [
        ("1", "Vá em Authentication → URL Configuration"),
        ("2", "Em Site URL, apague o valor atual"),
        ("3", "Cole: https://psicogest-one.vercel.app"),
        ("4", "Em Redirect URLs, clique em + Add URL"),
        ("5", "Cole: https://psicogest-one.vercel.app/**"),
        ("6", "Clique em Save"),
    ]
    for i, (num, texto) in enumerate(passos_url):
        y = 3.3 + i * 0.51
        rect(sl, 0.55, y, 0.32, 0.3, VERDE)
        txt(sl, num, 0.55, y, 0.32, 0.3,
            size=12, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(sl, texto, 0.97, y, 5.2, 0.42, size=12, color=TEXTO)

    # Box direito: confirmar e-mail
    rect(sl, 6.6, 1.6, 6.4, 5.5, RGBColor(0xFF, 0xF8, 0xE1))
    txt(sl, "📧  Ativar confirmação de e-mail", 6.75, 1.7, 6.2, 0.45,
        size=17, bold=True, color=RGBColor(0x7B, 0x5B, 0x00))
    txt(sl, "Sem isso, qualquer pessoa que se cadastrar entra direto no sistema sem confirmar o e-mail.",
        6.75, 2.22, 6.2, 0.65, size=13, color=TEXTO, wrap=True)
    txt(sl, "Como ativar:", 6.75, 2.95, 6.2, 0.4, size=14, bold=True, color=RGBColor(0x7B, 0x5B, 0x00))
    passos_email = [
        ("1", "Vá em Authentication → Providers"),
        ("2", "Clique em Email"),
        ("3", "Encontre o toggle Confirm email"),
        ("4", "Ligue o toggle (fica azul/verde)"),
        ("5", "Clique em Save"),
    ]
    for i, (num, texto) in enumerate(passos_email):
        y = 3.45 + i * 0.51
        rect(sl, 6.75, y, 0.32, 0.3, AMBAR)
        txt(sl, num, 6.75, y, 0.32, 0.3,
            size=12, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(sl, texto, 7.18, y, 5.8, 0.42, size=12, color=TEXTO)

    rodape(sl, "Esses dois ajustes resolvem: link de e-mail apontando para localhost + acesso sem confirmar e-mail")
    numero_slide(sl, 9, 13)


def slide_vercel_envvars(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, CINZA_F)
    cabecalho_faixa(sl, "Passo 7 — Adicionar variáveis de ambiente no Vercel",
                    "vercel.com  →  psicogest-one  →  Settings  →  Environment Variables", AZUL)

    txt(sl, "As variáveis de ambiente sao os segredos do seu sistema. O Vercel as guarda de forma segura.",
        0.4, 1.65, 12.5, 0.45, size=14, color=TEXTO)

    variaveis = [
        ("SUPABASE_SERVICE_ROLE_KEY",
         "Chave secreta do Supabase para acesso administrativo",
         "Supabase → Settings → API → service_role (a chave longa)"),
        ("STRIPE_SECRET_KEY",
         "Chave secreta do Stripe para processar pagamentos",
         "Stripe → Desenvolvedores → Chaves de API → Secret key (sk_live_...)"),
        ("STRIPE_PRICE_ID",
         "ID do preço criado no Stripe (o plano R$97/mês)",
         "Stripe → Produtos → clique no produto → copie o Price ID (price_...)"),
        ("STRIPE_WEBHOOK_SECRET",
         "Chave de verificação do webhook do Stripe",
         "Stripe → Desenvolvedores → Webhooks → clique no webhook → Signing secret (whsec_...)"),
        ("NEXT_PUBLIC_SITE_URL",
         "URL pública do seu sistema no Vercel",
         "https://psicogest-one.vercel.app"),
    ]

    for i, (nome, descricao, onde) in enumerate(variaveis):
        y = 2.25 + i * 0.97
        rect(sl, 0.4, y, 12.5, 0.9, BRANCO)
        rect(sl, 0.4, y, 0.08, 0.9, AZUL)
        rect(sl, 0.5, y + 0.05, 4.8, 0.38, RGBColor(0x1A, 0x1A, 0x2E))
        txt(sl, nome, 0.6, y + 0.07, 4.7, 0.3,
            size=12, bold=True, color=RGBColor(0x7D, 0xFF, 0xD8))
        txt(sl, descricao, 5.5, y + 0.05, 7.2, 0.35, size=12, bold=True, color=AZUL_ESC)
        txt(sl, f"Onde pegar: {onde}", 5.5, y + 0.47, 7.2, 0.37, size=11, color=CINZA_M)

    rodape(sl, "Após adicionar todas as variáveis → clique em Deployments → clique nos 3 pontos do último deploy → Redeploy")
    numero_slide(sl, 10, 13)


def slide_vercel_redeploy(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, CINZA_F)
    cabecalho_faixa(sl, "Passo 8 — Fazer Redeploy no Vercel",
                    "Necessário para as novas variáveis entrarem em funcionamento", AZUL)

    txt(sl, "Depois de salvar as variáveis no Vercel, o sistema precisa ser reiniciado para usá-las.",
        0.4, 1.65, 12.5, 0.45, size=15, color=TEXTO)

    # Passos visuais
    passos = [
        (AZUL,    "1", "Abra o Vercel",
         "vercel.com → faça login → abra o projeto psicogest-one"),
        (AZUL,    "2", "Vá em Deployments",
         "Clique na aba Deployments no menu superior"),
        (ROXO,    "3", "Abra o último deploy",
         "Clique nos 3 pontinhos (⋯) ao lado do deploy mais recente"),
        (ROXO,    "4", "Clique em Redeploy",
         "Selecione Redeploy no menu que aparece"),
        (VERDE,   "5", "Aguarde o build",
         "O Vercel vai reconstruir o sistema (leva 2 a 3 minutos)"),
        (VERDE,   "6", "Pronto!",
         "Quando aparecer Ready em verde, o sistema está no ar com as novas configurações"),
    ]

    for i, (cor, num, titulo, desc) in enumerate(passos):
        col = i % 3
        row = i // 3
        x = 0.4 + col * 4.3
        y = 2.25 + row * 2.3
        rect(sl, x, y, 3.9, 2.0, cor)
        rect(sl, x, y + 0.65, 3.9, 1.35, BRANCO)
        rect(sl, x, y, 0.55, 0.65, RGBColor(
            max(cor[0] - 30, 0), max(cor[1] - 30, 0), max(cor[2] - 30, 0)
        ))
        txt(sl, num, x, y, 0.55, 0.65,
            size=22, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(sl, titulo, x + 0.6, y + 0.12, 3.2, 0.45,
            size=15, bold=True, color=BRANCO)
        txt(sl, desc, x + 0.1, y + 0.72, 3.7, 1.18, size=12, color=TEXTO, wrap=True)

    rodape(sl, "Aguarde o status mudar para Ready (ícone verde) antes de testar o sistema")
    numero_slide(sl, 11, 13)


def slide_testar(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, CINZA_F)
    cabecalho_faixa(sl, "Passo 9 — Testar o sistema de pagamento",
                    "Use o cartão de teste do Stripe antes de cobrar de verdade", VERDE_ESC)

    # Cartão de teste
    rect(sl, 0.4, 1.6, 6.0, 3.2, RGBColor(0x1A, 0x1A, 0x2E))
    rect(sl, 0.4, 1.6, 6.0, 0.65, ROXO)
    txt(sl, "💳  Cartão de Teste do Stripe", 0.55, 1.67, 5.7, 0.5,
        size=16, bold=True, color=BRANCO)
    txt_multi(sl, [
        "Número:   4242 4242 4242 4242",
        "Validade: qualquer data futura (ex: 12/29)",
        "CVC:      qualquer 3 dígitos (ex: 123)",
        "Nome:     qualquer nome",
        "CEP:      qualquer CEP (ex: 01310-100)",
    ], 0.55, 2.35, 5.7, 2.2, size=14, color=RGBColor(0x7D, 0xFF, 0xD8))

    # Como testar
    rect(sl, 6.7, 1.6, 6.3, 3.2, BRANCO)
    txt(sl, "Como testar:", 6.85, 1.7, 6.1, 0.45, size=17, bold=True, color=VERDE_ESC)
    passos_t = [
        ("1", "Abra o sistema e crie uma conta nova"),
        ("2", "O trial expira simulando 14 dias"),
        ("3", "Clique em Assinar agora"),
        ("4", "Na página do Stripe, use o cartão acima"),
        ("5", "Conclua o pagamento"),
        ("6", "Verifique se o sistema liberou o acesso"),
    ]
    for i, (num, texto) in enumerate(passos_t):
        y = 2.18 + i * 0.46
        rect(sl, 6.85, y, 0.32, 0.3, VERDE)
        txt(sl, num, 6.85, y, 0.32, 0.3,
            size=12, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(sl, texto, 7.28, y, 5.6, 0.38, size=12, color=TEXTO)

    # Aviso produção
    rect(sl, 0.4, 5.0, 12.5, 2.1, RGBColor(0xFF, 0xF3, 0xCD))
    txt(sl, "⚠️  Quando estiver tudo funcionando no modo TESTE, troque para PRODUÇÃO:",
        0.55, 5.08, 12.2, 0.4, size=14, bold=True, color=RGBColor(0x7B, 0x5B, 0x00))
    txt_multi(sl, [
        "1. No Stripe, desative o Modo Teste (toggle no canto superior direito)",
        "2. Copie as novas chaves de PRODUÇÃO (sk_live_... e o novo webhook)",
        "3. Atualize STRIPE_SECRET_KEY e STRIPE_WEBHOOK_SECRET no Vercel",
        "4. Faça Redeploy   →   agora o sistema cobra de verdade!",
    ], 0.55, 5.55, 12.2, 1.45, size=13, color=RGBColor(0x7B, 0x5B, 0x00))

    numero_slide(sl, 12, 13)


def slide_resumo_final(prs):
    sl = slide_branco(prs)
    rect(sl, 0, 0, 13.333, 7.5, AZUL_ESC)
    rect(sl, 0, 0, 0.25, 7.5, VERDE)

    txt(sl, "Resumo — Checklist completo", 0.6, 0.3, 12, 0.7,
        size=28, bold=True, color=BRANCO)

    checks = [
        (ROXO,      "Stripe",   [
            "✅ Conta criada no stripe.com",
            "✅ Produto PsicoDoc Mensal criado (R$97/mês)",
            "✅ Price ID copiado (price_...)",
            "✅ Secret Key copiada (sk_live_...)",
            "✅ Webhook configurado com os 5 eventos",
            "✅ Webhook Signing Secret copiado (whsec_...)",
        ]),
        (VERDE,     "Supabase", [
            "✅ SQL da migração 004 executado",
            "✅ Apareceu Success no SQL Editor",
            "✅ Site URL atualizado para o domínio Vercel",
            "✅ Redirect URL adicionada (domínio/**)",
            "✅ Confirmação de e-mail ativada",
            "✅ Seu e-mail marcado como is_admin",
        ]),
        (AZUL,      "Vercel",   [
            "✅ SUPABASE_SERVICE_ROLE_KEY adicionada",
            "✅ STRIPE_SECRET_KEY adicionada",
            "✅ STRIPE_PRICE_ID adicionado",
            "✅ STRIPE_WEBHOOK_SECRET adicionado",
            "✅ NEXT_PUBLIC_SITE_URL adicionado",
            "✅ Redeploy feito e status Ready",
        ]),
    ]

    for i, (cor, titulo, items) in enumerate(checks):
        x = 0.5 + i * 4.25
        rect(sl, x, 1.2, 3.9, 0.55, cor)
        txt(sl, titulo, x, 1.2, 3.9, 0.55,
            size=17, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        rect(sl, x, 1.75, 3.9, 5.3, RGBColor(0x16, 0x24, 0x35))
        txt_multi(sl, items, x + 0.12, 1.85, 3.7, 5.1,
                  size=13, color=BRANCO)

    txt(sl, "🎉  Feito tudo isso, seu sistema já aceita pagamentos automáticos no cartão!",
        0.5, 7.05, 12.3, 0.38,
        size=14, bold=True, color=AMBAR, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# GERAÇÃO
# ════════════════════════════════════════════════════════════════

def gerar():
    prs = nova_apresentacao()

    slide_capa(prs)                  # 1
    slide_visao_geral(prs)           # 2
    slide_stripe_conta(prs)          # 3
    slide_stripe_produto(prs)        # 4
    slide_stripe_chaves(prs)         # 5
    slide_stripe_webhook(prs)        # 6
    slide_supabase_migracao(prs)     # 7
    slide_sql_codigo(prs)            # 8
    slide_supabase_url(prs)          # 9
    slide_vercel_envvars(prs)        # 10
    slide_vercel_redeploy(prs)       # 11
    slide_testar(prs)                # 12
    slide_resumo_final(prs)          # 13

    saida = "Guia_Pagamento_PsicoDoc.pptx"
    prs.save(saida)
    print(f"✅  Apresentação salva: {saida}")

if __name__ == "__main__":
    gerar()
