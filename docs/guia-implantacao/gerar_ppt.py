"""
Gerador do PPT de Implantação do PsicoGest
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import sys

# ── Paleta de cores ──────────────────────────────────────────────
AZUL       = RGBColor(0x4F, 0x7C, 0xAC)   # primária
VERDE      = RGBColor(0x6B, 0xAE, 0x8E)   # sucesso/passo concluído
AMBAR      = RGBColor(0xE8, 0xA8, 0x38)   # atenção
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_F    = RGBColor(0xF8, 0xF9, 0xFB)   # fundo claro
TEXTO      = RGBColor(0x1E, 0x2A, 0x38)   # texto escuro
CINZA_M    = RGBColor(0x64, 0x74, 0x8B)   # texto secundário
AZUL_ESC   = RGBColor(0x1E, 0x40, 0x6E)   # azul escuro (gradiente)
LARANJA    = RGBColor(0xD9, 0x55, 0x55)   # destaque negativo

W = Inches(13.333)   # largura 16:9
H = Inches(7.5)      # altura  16:9

# ── Helpers ──────────────────────────────────────────────────────

def nova_apresentacao():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def slide_branco(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # blank

def rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def txt(slide, text, left, top, width, height,
        size=18, bold=False, color=TEXTO, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def txt_multi(slide, lines, left, top, width, height,
              size=16, bold=False, color=TEXTO, align=PP_ALIGN.LEFT,
              line_space=None):
    """Cria textbox com múltiplos parágrafos."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox

def numero_slide(slide, n, total):
    txt(slide, f"{n} / {total}", 12.2, 7.1, 1.0, 0.35,
        size=11, color=CINZA_M, align=PP_ALIGN.RIGHT)

def header_barra(slide, cor=AZUL):
    """Barra colorida no topo."""
    rect(slide, 0, 0, 13.333, 0.08, cor)

def chip(slide, label, left, top, width=2.4, height=0.38,
         bg=AZUL, fg=BRANCO, size=12):
    r = rect(slide, left, top, width, height, bg)
    # arredondamento visual via ajuste de forma
    txt(slide, label, left + 0.1, top + 0.03, width - 0.2, height,
        size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════

def slide_capa(prs):
    s = slide_branco(prs)
    # fundo gradiente simulado (dois rects)
    rect(s, 0, 0, 13.333, 7.5, AZUL_ESC)
    rect(s, 0, 3.5, 13.333, 4.0, AZUL)

    # logo / nome
    txt(s, "PsicoGest", 0.8, 0.7, 6.0, 1.1,
        size=52, bold=True, color=BRANCO)
    txt(s, "Prontuário Eletrônico para Psicólogos",
        0.8, 1.75, 8.0, 0.7, size=22, color=RGBColor(0xC5, 0xDA, 0xF0))

    # linha divisória
    rect(s, 0.8, 2.55, 5.5, 0.05, VERDE)

    # título do guia
    txt(s, "Guia Completo de Implantação",
        0.8, 2.8, 9.0, 0.9, size=30, bold=True, color=BRANCO)
    txt(s, "Passo a passo para colocar o sistema no ar — sem precisar saber programar",
        0.8, 3.7, 10.0, 0.65, size=17, color=RGBColor(0xC5, 0xDA, 0xF0))

    # caixas de info
    chip(s, "⏱  Tempo estimado: 40 min", 0.8, 4.8, 3.6, 0.5,
         bg=RGBColor(0x1A, 0x35, 0x5C), fg=BRANCO, size=14)
    chip(s, "💰  Custo: R$ 0 (plano gratuito)", 4.6, 4.8, 3.8, 0.5,
         bg=RGBColor(0x1A, 0x35, 0x5C), fg=BRANCO, size=14)
    chip(s, "🔒  LGPD: Dados no Brasil", 8.6, 4.8, 3.5, 0.5,
         bg=RGBColor(0x1A, 0x35, 0x5C), fg=BRANCO, size=14)

    txt(s, "Versão 1.0 · Abril 2026", 0.8, 6.9, 5.0, 0.4,
        size=11, color=RGBColor(0x80, 0xA8, 0xCC))


def slide_visao_geral(prs, n, total):
    s = slide_branco(prs)
    header_barra(s)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    txt(s, "O que você vai precisar", 0.7, 0.3, 10.0, 0.75,
        size=32, bold=True, color=TEXTO)
    txt(s, "Nenhum conhecimento técnico necessário. Apenas estes itens:",
        0.7, 1.05, 10.0, 0.5, size=16, color=CINZA_M)

    # 3 colunas de pré-requisitos
    cols = [
        ("📧\nConta de e-mail\n(Gmail recomendado)",
         "Você vai criar contas em\n3 serviços gratuitos\ncom o mesmo e-mail."),
        ("🌐\nNavegador web\n(Chrome ou Edge)",
         "Todo o processo\nacontece no navegador.\nNada para instalar."),
        ("📁\nArquivos do sistema\n(pasta do projeto)",
         "A pasta 'Sistema para\npsicólogos' com todos\nos arquivos do PsicoGest."),
    ]
    for i, (titulo, desc) in enumerate(cols):
        x = 0.7 + i * 4.1
        rect(s, x, 1.7, 3.8, 2.5, BRANCO)
        txt_multi(s, titulo.split('\n'), x + 0.2, 1.85, 3.4, 1.0,
                  size=15, bold=True, color=TEXTO)
        txt_multi(s, desc.split('\n'), x + 0.2, 2.95, 3.4, 1.1,
                  size=13, color=CINZA_M)

    # Mapa dos passos
    txt(s, "Os 6 passos da implantação:", 0.7, 4.45, 10.0, 0.5,
        size=15, bold=True, color=TEXTO)

    passos = ["1\nGitHub", "2\nSupabase\n(banco)", "3\nSQL", "4\nStorage", "5\nVercel\n(deploy)", "6\nFinalizar"]
    cores = [AZUL, AZUL, VERDE, VERDE, AMBAR, VERDE]
    for i, (p, c) in enumerate(zip(passos, cores)):
        x = 0.7 + i * 2.1
        rect(s, x, 5.05, 1.8, 1.7, c)
        txt_multi(s, p.split('\n'), x + 0.05, 5.15, 1.7, 1.5,
                  size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        if i < 5:
            txt(s, "→", x + 1.82, 5.5, 0.3, 0.5,
                size=20, bold=True, color=CINZA_M)

    numero_slide(s, n, total)


def slide_github(prs, n, total):
    s = slide_branco(prs)
    header_barra(s)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "PASSO 1", 0.7, 0.25, 1.6, 0.4, bg=AZUL)
    txt(s, "Criar conta no GitHub e enviar o código",
        2.5, 0.25, 10.0, 0.55, size=24, bold=True, color=TEXTO)
    txt(s, "O GitHub guarda os arquivos do sistema. A Vercel (próximo passo) lê daqui para publicar.",
        0.7, 0.85, 12.0, 0.45, size=14, color=CINZA_M)

    passos = [
        ("1", "Acesse github.com e clique em 'Sign up'",
         "Use seu e-mail e crie uma senha. Confirme pelo e-mail recebido."),
        ("2", "Crie um novo repositório",
         "Clique em '+' → 'New repository'. Nome: psicogest. Marque 'Private'. Clique 'Create repository'."),
        ("3", "Faça upload dos arquivos",
         "Na página do repositório, clique 'uploading an existing file'. Arraste a pasta inteira do PsicoGest. Clique 'Commit changes'."),
        ("4", "Aguarde o upload concluir",
         "Pode demorar alguns minutos dependendo da sua internet. Ao terminar, os arquivos aparecem listados."),
    ]

    for i, (num, titulo, desc) in enumerate(passos):
        y = 1.45 + i * 1.45
        # círculo numerado
        rect(s, 0.7, y, 0.55, 0.55, AZUL)
        txt(s, num, 0.7, y, 0.55, 0.55,
            size=16, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        # conteúdo
        txt(s, titulo, 1.45, y, 11.0, 0.4,
            size=14, bold=True, color=TEXTO)
        txt(s, desc, 1.45, y + 0.42, 11.0, 0.75,
            size=12, color=CINZA_M)

    # dica
    rect(s, 0.7, 7.0, 11.7, 0.38, RGBColor(0xEF, 0xF6, 0xFF))
    txt(s, "💡  Não sabe usar Git? Tudo bem! O upload pela interface web do GitHub é suficiente.",
        0.85, 7.02, 11.4, 0.35, size=12, color=AZUL)

    numero_slide(s, n, total)


def slide_supabase_conta(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, VERDE)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "PASSO 2", 0.7, 0.25, 1.6, 0.4, bg=VERDE)
    txt(s, "Criar conta e projeto no Supabase (banco de dados)",
        2.5, 0.25, 10.0, 0.55, size=24, bold=True, color=TEXTO)
    txt(s, "O Supabase armazena todos os dados dos pacientes com segurança, no Brasil.",
        0.7, 0.85, 12.0, 0.45, size=14, color=CINZA_M)

    passos = [
        ("1", "Acesse supabase.com → 'Start your project'",
         "Faça login com sua conta do GitHub (clique 'Continue with GitHub'). Autorize o acesso."),
        ("2", "Clique em 'New project'",
         "Preencha: Nome do projeto = psicogest | Senha do banco = crie uma senha forte (anote!) | Região = South America (São Paulo)"),
        ("3", "Aguarde a criação do projeto",
         "Leva cerca de 2 minutos. Uma barra de progresso aparece. Não feche a aba!"),
        ("4", "Copie as credenciais",
         "No painel do projeto, clique em 'Project Settings' → 'API'. Copie o 'Project URL' e o 'anon public key'. Você vai precisar deles no Passo 6."),
    ]

    for i, (num, titulo, desc) in enumerate(passos):
        y = 1.45 + i * 1.45
        rect(s, 0.7, y, 0.55, 0.55, VERDE)
        txt(s, num, 0.7, y, 0.55, 0.55,
            size=16, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(s, titulo, 1.45, y, 11.0, 0.4,
            size=14, bold=True, color=TEXTO)
        txt(s, desc, 1.45, y + 0.42, 11.0, 0.75,
            size=12, color=CINZA_M)

    rect(s, 0.7, 7.0, 11.7, 0.38, RGBColor(0xF0, 0xFD, 0xF4))
    txt(s, "🔒  Região São Paulo = dados dos seus pacientes ficam no Brasil (exigência LGPD).",
        0.85, 7.02, 11.4, 0.35, size=12, color=VERDE)

    numero_slide(s, n, total)


def slide_sql(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, VERDE)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "PASSO 3", 0.7, 0.25, 1.6, 0.4, bg=VERDE)
    txt(s, "Criar as tabelas do banco de dados (SQL)",
        2.5, 0.25, 10.0, 0.55, size=24, bold=True, color=TEXTO)
    txt(s, "Você não precisa entender SQL. É só copiar e colar um arquivo pronto.",
        0.7, 0.85, 12.0, 0.45, size=14, color=CINZA_M)

    passos = [
        ("1", "No Supabase, abra o 'SQL Editor'",
         "No menu lateral esquerdo, clique no ícone de banco de dados → 'SQL Editor'. Uma tela de código vai aparecer."),
        ("2", "Abra o arquivo de migração",
         "No seu computador, abra a pasta do PsicoGest → supabase → migrations → '001_initial_schema.sql'. Abra com o Bloco de Notas."),
        ("3", "Copie todo o conteúdo",
         "Selecione tudo (Ctrl+A) e copie (Ctrl+C). Depois cole (Ctrl+V) na tela do SQL Editor do Supabase."),
        ("4", "Clique em 'Run' (botão verde)",
         "O Supabase vai criar todas as tabelas e configurações de segurança automaticamente. Aparecerá 'Success. No rows returned'."),
    ]

    for i, (num, titulo, desc) in enumerate(passos):
        y = 1.45 + i * 1.45
        rect(s, 0.7, y, 0.55, 0.55, VERDE)
        txt(s, num, 0.7, y, 0.55, 0.55,
            size=16, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(s, titulo, 1.45, y, 11.0, 0.4,
            size=14, bold=True, color=TEXTO)
        txt(s, desc, 1.45, y + 0.42, 11.0, 0.75,
            size=12, color=CINZA_M)

    rect(s, 0.7, 7.0, 11.7, 0.38, RGBColor(0xF0, 0xFD, 0xF4))
    txt(s, "✅  Após o 'Run', as tabelas de pacientes, consultas, prontuários e documentos estarão criadas.",
        0.85, 7.02, 11.4, 0.35, size=12, color=VERDE)

    numero_slide(s, n, total)


def slide_storage(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, VERDE)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "PASSO 4", 0.7, 0.25, 1.6, 0.4, bg=VERDE)
    txt(s, "Criar o espaço para arquivos (Storage)",
        2.5, 0.25, 10.0, 0.55, size=24, bold=True, color=TEXTO)
    txt(s, "Aqui ficam guardados os laudos, relatórios e documentos dos pacientes.",
        0.7, 0.85, 12.0, 0.45, size=14, color=CINZA_M)

    passos = [
        ("1", "No Supabase, clique em 'Storage' no menu lateral",
         "O ícone parece uma pasta/nuvem. A tela de armazenamento vai abrir."),
        ("2", "Clique em 'New bucket' (novo compartimento)",
         "Na tela que abrir, preencha: Nome = documents | Marque 'Public bucket' ✓ | Clique 'Save'."),
        ("3", "Confirme que o bucket foi criado",
         "Você verá 'documents' listado na tela de Storage. Isso é tudo!"),
    ]

    for i, (num, titulo, desc) in enumerate(passos):
        y = 1.5 + i * 1.6
        rect(s, 0.7, y, 0.55, 0.55, VERDE)
        txt(s, num, 0.7, y, 0.55, 0.55,
            size=16, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(s, titulo, 1.45, y, 11.0, 0.4,
            size=14, bold=True, color=TEXTO)
        txt(s, desc, 1.45, y + 0.42, 11.0, 0.75,
            size=12, color=CINZA_M)

    # resumo dos passos 2-4
    rect(s, 0.7, 6.35, 11.7, 0.55, RGBColor(0xE8, 0xF5, 0xE9))
    txt(s, "📋  Supabase concluído! Você já tem: banco de dados ✓   tabelas criadas ✓   storage configurado ✓",
        0.85, 6.42, 11.4, 0.42, size=13, bold=True, color=VERDE)

    numero_slide(s, n, total)


def slide_vercel(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, AMBAR)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "PASSO 5", 0.7, 0.25, 1.6, 0.4, bg=AMBAR)
    txt(s, "Publicar o sistema na Vercel (o site no ar)",
        2.5, 0.25, 10.0, 0.55, size=24, bold=True, color=TEXTO)
    txt(s, "A Vercel publica o PsicoGest na internet de forma gratuita. É o provedor recomendado para este sistema.",
        0.7, 0.85, 12.0, 0.45, size=14, color=CINZA_M)

    passos = [
        ("1", "Acesse vercel.com → 'Start Deploying'",
         "Faça login com sua conta do GitHub (botão 'Continue with GitHub'). Autorize o acesso quando solicitado."),
        ("2", "Importe o repositório do PsicoGest",
         "Clique 'Add New → Project'. Procure 'psicogest' na lista e clique 'Import'."),
        ("3", "Configure o projeto",
         "Em 'Framework Preset' selecione 'Next.js'. Deixe as outras opções como estão. NÃO clique em Deploy ainda."),
        ("4", "Adicione as variáveis de ambiente (IMPORTANTE)",
         "Clique em 'Environment Variables'. Você vai adicionar 2 variáveis — veja o próximo slide."),
    ]

    for i, (num, titulo, desc) in enumerate(passos):
        y = 1.45 + i * 1.45
        rect(s, 0.7, y, 0.55, 0.55, AMBAR)
        txt(s, num, 0.7, y, 0.55, 0.55,
            size=16, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(s, titulo, 1.45, y, 11.0, 0.4,
            size=14, bold=True, color=TEXTO)
        txt(s, desc, 1.45, y + 0.42, 11.0, 0.75,
            size=12, color=CINZA_M)

    rect(s, 0.7, 7.0, 11.7, 0.38, RGBColor(0xFF, 0xF8, 0xE8))
    txt(s, "⚠️  Não clique em Deploy antes de adicionar as variáveis de ambiente no Passo 6!",
        0.85, 7.02, 11.4, 0.35, size=12, bold=True, color=AMBAR)

    numero_slide(s, n, total)


def slide_variaveis(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, AMBAR)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "PASSO 6", 0.7, 0.25, 1.6, 0.4, bg=AMBAR)
    txt(s, "Configurar as variáveis de ambiente",
        2.5, 0.25, 10.0, 0.55, size=24, bold=True, color=TEXTO)
    txt(s, "São as 'senhas' que conectam a Vercel ao Supabase. Você obteve esses valores no Passo 2.",
        0.7, 0.85, 12.0, 0.45, size=14, color=CINZA_M)

    # tabela de variáveis
    rect(s, 0.7, 1.4, 11.7, 0.45, AZUL)
    txt(s, "Nome da Variável", 0.85, 1.47, 5.0, 0.35,
        size=13, bold=True, color=BRANCO)
    txt(s, "Onde encontrar no Supabase", 6.0, 1.47, 6.3, 0.35,
        size=13, bold=True, color=BRANCO)

    linhas = [
        ("NEXT_PUBLIC_SUPABASE_URL",
         "Project Settings → API → 'Project URL'\nEx: https://xyzabc.supabase.co"),
        ("NEXT_PUBLIC_SUPABASE_ANON_KEY",
         "Project Settings → API → 'anon public'\nChave longa começando com 'eyJh...'"),
    ]
    for i, (nome, onde) in enumerate(linhas):
        y = 1.85 + i * 1.2
        bg = BRANCO if i % 2 == 0 else RGBColor(0xF3, 0xF7, 0xFB)
        rect(s, 0.7, y, 11.7, 1.1, bg)
        rect(s, 0.7, y, 0.06, 1.1, AMBAR)  # borda esquerda colorida
        txt(s, nome, 0.95, y + 0.08, 5.0, 0.4,
            size=13, bold=True, color=TEXTO)
        txt_multi(s, onde.split('\n'), 6.0, y + 0.08, 6.3, 0.9,
                  size=12, color=CINZA_M)

    # instrução de deploy
    rect(s, 0.7, 4.35, 11.7, 1.5, RGBColor(0x1E, 0x2A, 0x38))
    txt(s, "Após adicionar as 2 variáveis:", 0.95, 4.5, 11.0, 0.4,
        size=14, bold=True, color=BRANCO)
    txt_multi(s, [
        "1.  Clique no botão 'Deploy' (azul, na parte inferior da página).",
        "2.  Aguarde cerca de 3 minutos enquanto o sistema é publicado.",
        "3.  Ao finalizar, a Vercel mostra confetes 🎉 e um link para o seu sistema!"
    ], 0.95, 4.95, 11.0, 0.8, size=13, color=RGBColor(0xC5, 0xDA, 0xF0))

    rect(s, 0.7, 5.95, 11.7, 0.7, RGBColor(0xF0, 0xFD, 0xF4))
    txt(s, "✅  Seu link ficará parecido com: https://psicogest-seunome.vercel.app",
        0.9, 6.05, 11.0, 0.45, size=14, bold=True, color=VERDE)

    numero_slide(s, n, total)


def slide_hostinger(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, AZUL)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "ALTERNATIVA", 0.7, 0.25, 2.1, 0.4, bg=AZUL)
    txt(s, "Deploy na Hostinger (plano pago)",
        3.0, 0.25, 10.0, 0.55, size=24, bold=True, color=TEXTO)

    # comparativo
    txt(s, "Por que a Vercel é recomendada?", 0.7, 0.9, 6.0, 0.45,
        size=16, bold=True, color=TEXTO)
    txt(s, "Quando usar a Hostinger?", 7.0, 0.9, 6.0, 0.45,
        size=16, bold=True, color=TEXTO)

    vercel_pts = [
        "✅  Gratuita para projetos assim",
        "✅  Feita para Next.js (mesma empresa)",
        "✅  Deploy automático a cada atualização",
        "✅  HTTPS automático e CDN global",
        "✅  Zero configuração de servidor",
    ]
    host_pts = [
        "✅  Você já tem a hospedagem paga",
        "✅  Quer domínio e e-mail profissional juntos",
        "✅  Prefere centralizar tudo num só lugar",
        "⚠️  Requer plano Business ou VPS (Node.js)",
        "⚠️  Configuração mais técnica necessária",
    ]

    rect(s, 0.7, 1.45, 5.8, 3.5, BRANCO)
    for i, pt in enumerate(vercel_pts):
        txt(s, pt, 0.9, 1.6 + i * 0.62, 5.5, 0.5, size=13, color=TEXTO)

    rect(s, 7.0, 1.45, 5.8, 3.5, BRANCO)
    for i, pt in enumerate(host_pts):
        c = TEXTO if "✅" in pt else AMBAR
        txt(s, pt, 7.2, 1.6 + i * 0.62, 5.5, 0.5, size=13, color=c)

    # instruções hostinger
    txt(s, "Se optar pela Hostinger:", 0.7, 5.1, 11.7, 0.4,
        size=15, bold=True, color=TEXTO)
    inst = [
        "1. Contrate o plano 'Business' ou superior (tem suporte a Node.js)",
        "2. No painel hPanel → Node.js → 'Create application' → Version 20 → App root: /",
        "3. Conecte ao GitHub ou faça upload via FTP (usando o FileZilla)",
        "4. Adicione as variáveis de ambiente no painel → Execute 'npm run build && npm start'",
    ]
    for i, inst_txt in enumerate(inst):
        txt(s, inst_txt, 0.7, 5.6 + i * 0.42, 12.0, 0.38, size=12, color=CINZA_M)

    numero_slide(s, n, total)


def slide_pronto(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, VERDE)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "CONCLUÍDO ✓", 0.7, 0.25, 2.2, 0.4, bg=VERDE)
    txt(s, "Sistema no ar! Primeiros passos",
        3.1, 0.25, 9.5, 0.55, size=24, bold=True, color=TEXTO)

    # Checklist final
    txt(s, "Checklist de verificação:", 0.7, 0.9, 8.0, 0.45,
        size=16, bold=True, color=TEXTO)

    checklist = [
        "Abra o link da Vercel no navegador — a landing page do PsicoGest deve aparecer",
        "Clique em 'Criar conta grátis' e cadastre-se como psicólogo",
        "Complete o onboarding (3 passos: perfil → paciente → agenda)",
        "Cadastre um paciente de teste e registre uma sessão",
        "Verifique se o upload de documento funciona na aba 'Documentos'",
        "Acesse pelo celular e confirme que a interface está responsiva",
    ]
    for i, item in enumerate(checklist):
        y = 1.45 + i * 0.62
        rect(s, 0.7, y + 0.06, 0.38, 0.38, VERDE)
        txt(s, "✓", 0.7, y + 0.04, 0.38, 0.38,
            size=14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(s, item, 1.2, y, 11.5, 0.5, size=13, color=TEXTO)

    # Próximos passos
    rect(s, 0.7, 5.4, 11.7, 1.7, RGBColor(0xEF, 0xF6, 0xFF))
    txt(s, "Próximos passos após validar com a Lilian:",
        0.9, 5.5, 10.0, 0.4, size=14, bold=True, color=AZUL)
    txt_multi(s, [
        "💰  Integrar pagamento (Stripe ou Asaas) para cobrar a assinatura mensal — Fase 4 do projeto",
        "📲  Configurar domínio personalizado (ex: psicogest.com.br) na Vercel",
        "📧  Configurar e-mails transacionais (boas-vindas, lembretes)",
    ], 0.9, 5.95, 11.3, 1.0, size=12, color=CINZA_M)

    numero_slide(s, n, total)


def slide_dicas(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, AZUL)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    txt(s, "Dúvidas frequentes e solução de problemas",
        0.7, 0.3, 11.5, 0.65, size=28, bold=True, color=TEXTO)

    faq = [
        ("❓ O deploy falhou na Vercel. O que fazer?",
         "Verifique se as 2 variáveis de ambiente foram adicionadas corretamente. Nomes e valores sem espaços extras."),
        ("❓ Cadastrei paciente mas não aparece nada.",
         "Confirme que o SQL foi executado no Supabase. Abra 'Table Editor' e veja se a tabela 'patients' existe."),
        ("❓ Upload de documento não funciona.",
         "Verifique se o bucket 'documents' foi criado no Supabase Storage (Passo 4). O nome deve ser exatamente 'documents'."),
        ("❓ Quero um domínio .com.br no lugar do .vercel.app.",
         "Compre o domínio na Hostinger (ou outro registrador). Na Vercel, vá em Settings → Domains e adicione o domínio comprado."),
        ("❓ Posso mudar a logo e o nome do sistema?",
         "Sim. Edite o arquivo app/(public)/page.tsx para mudar a landing page e components/app-shell.tsx para o logo interno."),
    ]

    for i, (pergunta, resposta) in enumerate(faq):
        y = 1.05 + i * 1.2
        rect(s, 0.7, y, 11.7, 1.1, BRANCO)
        rect(s, 0.7, y, 0.06, 1.1, AZUL)
        txt(s, pergunta, 0.9, y + 0.06, 11.3, 0.38,
            size=13, bold=True, color=TEXTO)
        txt(s, resposta, 0.9, y + 0.5, 11.3, 0.5,
            size=12, color=CINZA_M)

    numero_slide(s, n, total)


def slide_encerramento(prs):
    s = slide_branco(prs)
    rect(s, 0, 0, 13.333, 7.5, AZUL_ESC)
    rect(s, 0, 4.0, 13.333, 3.5, AZUL)

    txt(s, "PsicoGest", 0.8, 0.8, 8.0, 1.0,
        size=52, bold=True, color=BRANCO)
    rect(s, 0.8, 1.85, 4.0, 0.06, VERDE)

    txt(s, "Seu consultório organizado,\nsua mente tranquila.",
        0.8, 2.0, 8.0, 1.3, size=28, bold=True, color=BRANCO)

    txt(s, "Sistema implantado com sucesso? Compartilhe com outros psicólogos! 🎉",
        0.8, 4.2, 11.5, 0.6, size=18, color=RGBColor(0xC5, 0xDA, 0xF0))

    chip(s, "suporte: docs/PROGRESSO.md",
         0.8, 5.2, 4.5, 0.45, bg=RGBColor(0x1A, 0x35, 0x5C), fg=BRANCO, size=13)
    chip(s, "vercel.com  ·  supabase.com",
         5.5, 5.2, 4.0, 0.45, bg=RGBColor(0x1A, 0x35, 0x5C), fg=BRANCO, size=13)

    txt(s, "© 2026 PsicoGest · Versão 1.0",
        0.8, 6.9, 6.0, 0.4, size=11, color=RGBColor(0x80, 0xA8, 0xCC))


# ════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════

def main():
    prs = nova_apresentacao()
    total = 10  # total de slides com número

    slide_capa(prs)                      # slide 1 (sem número)
    slide_visao_geral(prs, 1, total)     # slide 2
    slide_github(prs, 2, total)          # slide 3
    slide_supabase_conta(prs, 3, total)  # slide 4
    slide_sql(prs, 4, total)             # slide 5
    slide_storage(prs, 5, total)         # slide 6
    slide_vercel(prs, 6, total)          # slide 7
    slide_variaveis(prs, 7, total)       # slide 8
    slide_hostinger(prs, 8, total)       # slide 9
    slide_pronto(prs, 9, total)          # slide 10
    slide_dicas(prs, 10, total)          # slide 11
    slide_encerramento(prs)              # slide 12 (sem número)

    out = "/home/ubuntu/Área de Trabalho/Sistema para psicólogos/docs/guia-implantacao/PsicoGest_Guia_Implantacao.pptx"
    prs.save(out)
    print(f"✅  Apresentação salva em:\n    {out}")
    print(f"    Slides: 12")


if __name__ == "__main__":
    main()
