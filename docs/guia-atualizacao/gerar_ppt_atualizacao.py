"""
Gerador do PPT de Atualização do PsicoDoc — Versão 2.0
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paleta de cores ──────────────────────────────────────────────
AZUL       = RGBColor(0x4F, 0x7C, 0xAC)
VERDE      = RGBColor(0x6B, 0xAE, 0x8E)
AMBAR      = RGBColor(0xE8, 0xA8, 0x38)
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_F    = RGBColor(0xF8, 0xF9, 0xFB)
TEXTO      = RGBColor(0x1E, 0x2A, 0x38)
CINZA_M    = RGBColor(0x64, 0x74, 0x8B)
AZUL_ESC   = RGBColor(0x1E, 0x40, 0x6E)
ROXO       = RGBColor(0x7B, 0x68, 0xC8)
VERMELHO   = RGBColor(0xD9, 0x55, 0x55)

W = Inches(13.333)
H = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────

def nova_apresentacao():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def slide_branco(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, left, top, width, height, fill_rgb):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def txt(slide, text, left, top, width, height,
        size=18, bold=False, color=TEXTO, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def txt_multi(slide, lines, left, top, width, height,
              size=16, bold=False, color=TEXTO, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox

def numero_slide(slide, n, total):
    txt(slide, f"{n} / {total}", 12.2, 7.1, 1.0, 0.35,
        size=11, color=CINZA_M, align=PP_ALIGN.RIGHT)

def header_barra(slide, cor=AZUL):
    rect(slide, 0, 0, 13.333, 0.08, cor)

def chip(slide, label, left, top, width=2.4, height=0.38,
         bg=AZUL, fg=BRANCO, size=12):
    rect(slide, left, top, width, height, bg)
    txt(slide, label, left + 0.1, top + 0.03, width - 0.2, height,
        size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)

def caixa_passo(slide, numero, titulo, subtitulo, left, top, cor=AZUL):
    rect(slide, left, top, 0.55, 0.55, cor)
    txt(slide, numero, left, top, 0.55, 0.55,
        size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    txt(slide, titulo, left + 0.65, top, 8.5, 0.3,
        size=14, bold=True, color=TEXTO)
    txt(slide, subtitulo, left + 0.65, top + 0.28, 8.5, 0.35,
        size=12, color=CINZA_M)

# ════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════

def slide_capa(prs):
    s = slide_branco(prs)
    rect(s, 0, 0, 13.333, 7.5, AZUL_ESC)
    rect(s, 0, 3.8, 13.333, 3.7, AZUL)

    txt(s, "PsicoDoc", 0.8, 0.7, 6.0, 1.0,
        size=52, bold=True, color=BRANCO)
    txt(s, "Prontuário Eletrônico para Psicólogos",
        0.8, 1.72, 9.0, 0.65, size=20, color=RGBColor(0xC5, 0xDA, 0xF0))

    rect(s, 0.8, 2.5, 5.5, 0.05, VERDE)

    txt(s, "Guia de Atualização — Versão 2.0",
        0.8, 2.75, 10.0, 0.9, size=30, bold=True, color=BRANCO)
    txt(s, "Passo a passo para atualizar o sistema com as 7 novas funcionalidades",
        0.8, 3.62, 10.5, 0.55, size=16, color=RGBColor(0xC5, 0xDA, 0xF0))

    chip(s, "⏱  Tempo estimado: 15 min", 0.8, 4.75, 3.6, 0.5,
         bg=RGBColor(0x1A, 0x35, 0x5C), fg=BRANCO, size=14)
    chip(s, "🗄️  1 SQL no Supabase", 4.6, 4.75, 3.1, 0.5,
         bg=RGBColor(0x1A, 0x35, 0x5C), fg=BRANCO, size=14)
    chip(s, "🚀  Deploy automático", 7.9, 4.75, 3.0, 0.5,
         bg=RGBColor(0x1A, 0x35, 0x5C), fg=BRANCO, size=14)

    chip(s, "✅  Nenhum conhecimento técnico necessário", 0.8, 5.5, 5.8, 0.5,
         bg=VERDE, fg=BRANCO, size=14)

    txt(s, "Versão 2.0 · Abril 2026", 0.8, 6.9, 5.0, 0.4,
        size=11, color=RGBColor(0x80, 0xA8, 0xCC))


def slide_novidades(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, VERDE)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    txt(s, "O que há de novo na Versão 2.0", 0.7, 0.25, 11.0, 0.7,
        size=28, bold=True, color=TEXTO)
    txt(s, "7 melhorias implementadas com base no feedback da psicóloga piloto",
        0.7, 0.95, 11.0, 0.45, size=14, color=CINZA_M)

    novidades = [
        (AZUL,  "1", "Campo de data de nascimento melhorado",
                      "Seleção por Dia / Mês / Ano — sem confusão de formatos"),
        (AZUL,  "2", "Status do paciente atualizado",
                      "Agora: Acompanhamento · Em Avaliação · Alta/Desligado"),
        (VERDE, "3", "Agendamentos recorrentes",
                      "Crie sessões semanais, quinzenais ou mensais de uma só vez"),
        (VERDE, "4", "Lembretes por WhatsApp",
                      "Mensagem automática 1 dia antes de cada sessão"),
        (AMBAR, "5", "Feliz aniversário automático",
                      "Mensagem personalizada enviada no dia do aniversário do paciente"),
        (ROXO,  "6", "Módulo de Contratos Terapêuticos",
                      "Crie e envie contratos para assinatura digital — sem papel"),
        (ROXO,  "7", "Aba de Laudo no prontuário",
                      "Formulário estruturado para laudos neuropsicológicos e psicológicos"),
    ]

    for i, (cor, num, titulo, desc) in enumerate(novidades):
        row = i % 4
        col = i // 4
        x = 0.7 + col * 6.55
        y = 1.55 + row * 1.35

        rect(s, x, y, 6.0, 1.15, BRANCO)
        rect(s, x, y, 0.45, 1.15, cor)
        txt(s, num, x, y + 0.3, 0.45, 0.6,
            size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(s, titulo, x + 0.55, y + 0.08, 5.35, 0.4,
            size=13, bold=True, color=TEXTO)
        txt(s, desc, x + 0.55, y + 0.48, 5.35, 0.55,
            size=11, color=CINZA_M)

    numero_slide(s, n, total)


def slide_visao_geral(prs, n, total):
    s = slide_branco(prs)
    header_barra(s)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    txt(s, "Visão geral da atualização", 0.7, 0.25, 11.0, 0.7,
        size=28, bold=True, color=TEXTO)
    txt(s, "Apenas 3 etapas principais — todo o processo leva cerca de 15 minutos",
        0.7, 0.95, 11.0, 0.45, size=14, color=CINZA_M)

    # Etapas principais
    etapas = [
        (AZUL,  "ETAPA 1\n\n🗄️\n\nBanco de\nDados\n(Supabase)",
                 "Executar o SQL\nde atualização"),
        (VERDE, "ETAPA 2\n\n📁\n\nCódigo\nAtualizado\n(GitHub)",
                 "Enviar os arquivos\nnovo para o GitHub"),
        (AMBAR, "ETAPA 3\n\n🚀\n\nDeploy\nAutomático\n(Vercel)",
                 "Verificar que o\nsistema subiu"),
    ]

    for i, (cor, titulo, desc) in enumerate(etapas):
        x = 1.2 + i * 3.9
        rect(s, x, 1.6, 3.3, 3.8, cor)
        txt_multi(s, titulo.split('\n'), x + 0.15, 1.75, 3.0, 3.3,
                  size=16, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        if i < 2:
            txt(s, "→", x + 3.35, 3.1, 0.5, 0.6,
                size=28, bold=True, color=CINZA_M)

    # Notas de rodapé
    notas = [
        "⚠️  A ETAPA 1 (SQL) é obrigatória: sem ela, o sistema não funcionará corretamente.",
        "✅  Se o sistema já estava conectado ao GitHub + Vercel, a ETAPA 3 acontece automaticamente.",
    ]
    for i, nota in enumerate(notas):
        cor_bg = AMBAR if i == 0 else VERDE
        rect(s, 0.7, 5.65 + i * 0.72, 11.9, 0.6, cor_bg)
        txt(s, nota, 0.9, 5.7 + i * 0.72, 11.5, 0.5,
            size=12, bold=False, color=BRANCO)

    numero_slide(s, n, total)


def slide_etapa1_sql(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, AZUL)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "ETAPA 1", 0.7, 0.25, 1.6, 0.4, bg=AZUL)
    txt(s, "Atualizar o banco de dados no Supabase",
        2.45, 0.25, 10.5, 0.55, size=24, bold=True, color=TEXTO)
    txt(s, "Este passo adiciona as novas tabelas e colunas necessárias para as 7 melhorias.",
        0.7, 0.85, 12.0, 0.4, size=13, color=CINZA_M)

    passos = [
        ("1", "Acesse o Supabase",
              "Abra supabase.com no navegador e faça login com sua conta."),
        ("2", "Abra seu projeto",
              "Clique no projeto do PsicoDoc (ex.: psicogest-producao)."),
        ("3", "Vá em SQL Editor",
              "No menu lateral esquerdo, clique em 'SQL Editor' (ícone de banco de dados)."),
        ("4", "Clique em 'New query'",
              "Botão no canto superior esquerdo da tela do SQL Editor."),
        ("5", "Cole o SQL de atualização",
              "Copie TODO o conteúdo do arquivo  supabase/migrations/002_features.sql  e cole na caixa de texto."),
        ("6", "Execute o SQL",
              "Clique no botão verde 'Run' (ou pressione Ctrl+Enter). Aguarde a mensagem 'Success'."),
    ]

    for i, (num, titulo, desc) in enumerate(passos):
        row = i % 3
        col = i // 3
        x = 0.7 + col * 6.4
        y = 1.4 + row * 1.85

        rect(s, x, y, 5.9, 1.65, BRANCO)
        rect(s, x, y, 0.5, 1.65, AZUL)
        txt(s, num, x, y + 0.5, 0.5, 0.65,
            size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(s, titulo, x + 0.6, y + 0.12, 5.1, 0.4,
            size=13, bold=True, color=TEXTO)
        txt(s, desc, x + 0.6, y + 0.52, 5.1, 0.9,
            size=11, color=CINZA_M)

    rect(s, 0.7, 7.0, 11.9, 0.35, RGBColor(0xFF, 0xF3, 0xCD))
    txt(s, "⚠️  Se aparecer erro 'already exists', ignore — significa que a tabela já foi criada anteriormente. Continue.",
        0.9, 7.02, 11.5, 0.3, size=11, color=RGBColor(0x7A, 0x5C, 0x00))

    numero_slide(s, n, total)


def slide_etapa2_github(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, VERDE)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "ETAPA 2", 0.7, 0.25, 1.6, 0.4, bg=VERDE)
    txt(s, "Enviar o código atualizado para o GitHub",
        2.45, 0.25, 10.5, 0.55, size=24, bold=True, color=TEXTO)
    txt(s, "O GitHub guarda todos os arquivos do sistema. Precisamos enviar a versão nova.",
        0.7, 0.85, 12.0, 0.4, size=13, color=CINZA_M)

    # Dois caminhos
    txt(s, "Escolha UMA das opções abaixo:", 0.7, 1.35, 10.0, 0.45,
        size=14, bold=True, color=TEXTO)

    # Opção A
    rect(s, 0.7, 1.85, 5.8, 4.8, BRANCO)
    rect(s, 0.7, 1.85, 5.8, 0.45, AZUL)
    txt(s, "OPÇÃO A — GitHub Desktop (recomendado)", 0.9, 1.9, 5.4, 0.4,
        size=12, bold=True, color=BRANCO)

    opcao_a = [
        "1. Abra o GitHub Desktop no computador",
        "2. Selecione o repositório do PsicoDoc",
        "3. Ele mostrará automaticamente todos\n    os arquivos alterados (em verde)",
        "4. No campo 'Summary' escreva:\n    Atualização versão 2.0",
        "5. Clique em 'Commit to main'",
        "6. Clique em 'Push origin'",
        "✅  Pronto! Código enviado.",
    ]
    for i, linha in enumerate(opcao_a):
        cor = VERDE if linha.startswith("✅") else TEXTO
        bold = linha.startswith("✅")
        txt_multi(s, linha.split('\n'), 0.9, 2.4 + i * 0.55, 5.3, 0.5,
                  size=12, bold=bold, color=cor)

    # Opção B
    rect(s, 6.85, 1.85, 5.8, 4.8, BRANCO)
    rect(s, 6.85, 1.85, 5.8, 0.45, CINZA_M)
    txt(s, "OPÇÃO B — Pelo site do GitHub", 7.05, 1.9, 5.4, 0.4,
        size=12, bold=True, color=BRANCO)

    opcao_b = [
        "1. Acesse github.com e abra seu repositório",
        "2. Clique em 'Add file' > 'Upload files'",
        "3. Arraste TODOS os arquivos/pastas\n    atualizados para a área de upload",
        "4. No campo 'Commit changes'\n    escreva: Atualização versão 2.0",
        "5. Clique em 'Commit changes'",
        "✅  Pronto! Código enviado.",
    ]
    for i, linha in enumerate(opcao_b):
        cor = VERDE if linha.startswith("✅") else TEXTO
        bold = linha.startswith("✅")
        txt_multi(s, linha.split('\n'), 7.05, 2.4 + i * 0.55, 5.3, 0.5,
                  size=12, bold=bold, color=cor)

    # Separador entre opções
    txt(s, "OU", 6.45, 3.8, 0.5, 0.5,
        size=16, bold=True, color=CINZA_M, align=PP_ALIGN.CENTER)

    numero_slide(s, n, total)


def slide_etapa3_vercel(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, AMBAR)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "ETAPA 3", 0.7, 0.25, 1.6, 0.4, bg=AMBAR)
    txt(s, "Verificar o deploy automático na Vercel",
        2.45, 0.25, 10.5, 0.55, size=24, bold=True, color=TEXTO)
    txt(s, "Após o envio para o GitHub, a Vercel detecta automaticamente e publica a nova versão.",
        0.7, 0.85, 12.0, 0.4, size=13, color=CINZA_M)

    # Passos de verificação
    txt(s, "Como verificar:", 0.7, 1.35, 6.0, 0.45,
        size=16, bold=True, color=TEXTO)

    passos_v = [
        ("1", AMBAR, "Acesse vercel.com",
              "Faça login com a mesma conta que você usou para fazer o deploy."),
        ("2", AMBAR, "Clique no projeto PsicoDoc",
              "Na lista de projetos, clique no card do PsicoDoc."),
        ("3", AMBAR, "Acompanhe o deploy",
              "Você verá um card 'Deployments'. O mais recente mostrará 'Building...' e depois 'Ready'."),
        ("4", VERDE, "Acesse o sistema",
              "Clique no link do seu domínio (ex.: psicodoc.vercel.app) para confirmar que está funcionando."),
    ]

    for i, (num, cor, titulo, desc) in enumerate(passos_v):
        y = 1.95 + i * 1.2
        rect(s, 0.7, y, 11.9, 1.05, BRANCO)
        rect(s, 0.7, y, 0.5, 1.05, cor)
        txt(s, num, 0.7, y + 0.25, 0.5, 0.55,
            size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(s, titulo, 1.3, y + 0.08, 10.1, 0.4,
            size=14, bold=True, color=TEXTO)
        txt(s, desc, 1.3, y + 0.5, 10.1, 0.45,
            size=12, color=CINZA_M)

    # Dica de tempo
    rect(s, 0.7, 6.85, 11.9, 0.5, RGBColor(0xE8, 0xF4, 0xED))
    txt(s, "⏱  Normalmente o deploy leva entre 1 e 3 minutos. Se demorar mais de 10 minutos, atualize a página.",
        0.9, 6.9, 11.5, 0.4, size=12, color=RGBColor(0x1A, 0x5C, 0x35))

    numero_slide(s, n, total)


def slide_testar(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, ROXO)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "VERIFICAÇÃO", 0.7, 0.25, 2.1, 0.4, bg=ROXO)
    txt(s, "Testar as novas funcionalidades",
        2.95, 0.25, 10.0, 0.55, size=24, bold=True, color=TEXTO)
    txt(s, "Após o deploy, faça um teste rápido para confirmar que tudo está funcionando:",
        0.7, 0.85, 12.0, 0.4, size=13, color=CINZA_M)

    testes = [
        (AZUL,  "Prontuário do paciente",
                 "Abra um paciente e confirme que a aba 'Laudo' aparece (5ª aba)"),
        (AZUL,  "Status dos pacientes",
                 "Na lista de pacientes, verifique se os status mostram 'Acompanhamento', 'Em Avaliação' ou 'Alta'"),
        (VERDE, "Agenda — sessão recorrente",
                 "Na agenda, clique em + Novo agendamento e veja a opção 'Repetir' (semanal/quinzenal/mensal)"),
        (VERDE, "Contratos",
                 "No menu lateral, clique em 'Contratos' e tente criar um novo contrato para um paciente"),
        (AMBAR, "Link de assinatura",
                 "Após criar um contrato, copie o link gerado, abra numa aba anônima e veja a tela de assinatura"),
        (ROXO,  "Configurações — WhatsApp",
                 "Acesse Configurações e role até o final para ver a seção de integração com WhatsApp"),
    ]

    for i, (cor, titulo, desc) in enumerate(testes):
        row = i % 3
        col = i // 3
        x = 0.7 + col * 6.4
        y = 1.4 + row * 1.85

        rect(s, x, y, 5.9, 1.6, BRANCO)
        rect(s, x, y, 0.45, 1.6, cor)

        # Checkbox
        rect(s, x + 0.1, y + 0.1, 0.25, 0.25, CINZA_F)
        txt(s, "☐", x + 0.1, y + 0.08, 0.3, 0.3, size=11, color=CINZA_M)

        txt(s, titulo, x + 0.6, y + 0.1, 5.1, 0.4,
            size=13, bold=True, color=TEXTO)
        txt(s, desc, x + 0.6, y + 0.52, 5.1, 0.9,
            size=11, color=CINZA_M)

    numero_slide(s, n, total)


def slide_whatsapp_config(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, VERDE)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "OPCIONAL", 0.7, 0.25, 1.8, 0.4, bg=CINZA_M)
    txt(s, "Configurar WhatsApp automático",
        2.65, 0.25, 10.5, 0.55, size=24, bold=True, color=TEXTO)
    txt(s, "Para ativar os lembretes e mensagens de aniversário, siga as instruções abaixo:",
        0.7, 0.85, 12.0, 0.4, size=13, color=CINZA_M)

    # Coluna esquerda: Evolution API
    rect(s, 0.7, 1.4, 5.9, 5.4, BRANCO)
    rect(s, 0.7, 1.4, 5.9, 0.45, VERDE)
    txt(s, "Evolution API (gratuito — self-hosted)", 0.9, 1.45, 5.5, 0.38,
        size=13, bold=True, color=BRANCO)

    ev_passos = [
        "1. Acesse railway.app e crie conta",
        "2. Clique 'New Project' > busque 'Evolution API'",
        "3. Deploy com os padrões sugeridos",
        "4. Após subir, copie a URL do projeto",
        "5. Acesse URL/instance/create e crie uma instância",
        "6. Escaneie o QR Code com o WhatsApp do consultório",
        "7. No PsicoDoc: Configurações > WhatsApp",
        "8. Cole: URL, nome da instância e token (apikey)",
        "9. Ative os lembretes e salve",
    ]
    for i, p in enumerate(ev_passos):
        txt(s, p, 0.9, 1.95 + i * 0.52, 5.5, 0.48, size=11, color=TEXTO)

    # Coluna direita: Z-API
    rect(s, 6.85, 1.4, 5.9, 5.4, BRANCO)
    rect(s, 6.85, 1.4, 5.9, 0.45, AZUL)
    txt(s, "Z-API (pago — mais simples)", 7.05, 1.45, 5.5, 0.38,
        size=13, bold=True, color=BRANCO)

    zapi_passos = [
        "1. Acesse z-api.io e crie uma conta",
        "2. Crie uma nova instância",
        "3. Escaneie o QR Code com seu WhatsApp",
        "4. Copie o ID da instância e o Token",
        "5. No PsicoDoc: Configurações > WhatsApp",
        "6. Selecione 'Z-API' como provedor",
        "7. Cole: URL (https://api.z-api.io),",
        "    ID da instância e Token",
        "8. Ative os lembretes e salve",
    ]
    for i, p in enumerate(zapi_passos):
        txt(s, p, 7.05, 1.95 + i * 0.52, 5.5, 0.48, size=11, color=TEXTO)

    # Nota
    rect(s, 0.7, 7.0, 11.9, 0.38, RGBColor(0xFF, 0xF3, 0xCD))
    txt(s, "⚠️  Os lembretes precisam ser disparados manualmente ou via cron job. Veja o slide de automação na documentação técnica.",
        0.9, 7.03, 11.5, 0.3, size=10, color=RGBColor(0x7A, 0x5C, 0x00))

    numero_slide(s, n, total)


def slide_problemas(prs, n, total):
    s = slide_branco(prs)
    header_barra(s, VERMELHO)
    rect(s, 0, 0.08, 13.333, 7.42, CINZA_F)

    chip(s, "SUPORTE", 0.7, 0.25, 1.6, 0.4, bg=VERMELHO)
    txt(s, "Problemas comuns e como resolver",
        2.45, 0.25, 10.5, 0.55, size=24, bold=True, color=TEXTO)

    problemas = [
        (AMBAR, "O SQL deu erro 'already exists'",
                 "Sem problema! Significa que essa parte já foi criada antes. Você pode ignorar esse erro e continuar.",
                 "✅  Solução: Ignore o erro e execute o restante do script."),
        (AMBAR, "O deploy na Vercel ficou travado em 'Building'",
                 "Pode acontecer em caso de erro no código. Acesse vercel.com > seu projeto > Deployments > clique no deploy com erro.",
                 "✅  Solução: Leia os logs de erro e entre em contato com o suporte."),
        (VERMELHO, "O site abriu mas está com erro 500",
                 "Geralmente significa que as variáveis de ambiente estão faltando ou o SQL não foi executado.",
                 "✅  Solução: Confira se a ETAPA 1 (SQL) foi concluída e verifique as variáveis na Vercel."),
        (VERMELHO, "Os pacientes sumiram ou o status está errado",
                 "O SQL de migração atualiza os status antigos (active → acompanhamento). Execute o SQL novamente.",
                 "✅  Solução: Execute novamente o arquivo 002_features.sql no Supabase."),
    ]

    for i, (cor, titulo, desc, solucao) in enumerate(problemas):
        row = i % 2
        col = i // 2
        x = 0.7 + col * 6.4
        y = 1.35 + row * 2.7

        rect(s, x, y, 5.9, 2.45, BRANCO)
        rect(s, x, y, 5.9, 0.38, cor)
        txt(s, titulo, x + 0.15, y + 0.07, 5.6, 0.3,
            size=13, bold=True, color=BRANCO)
        txt(s, desc, x + 0.15, y + 0.48, 5.6, 0.9,
            size=11, color=CINZA_M)
        txt(s, solucao, x + 0.15, y + 1.62, 5.6, 0.55,
            size=11, bold=True, color=VERDE)

    numero_slide(s, n, total)


def slide_final(prs, n, total):
    s = slide_branco(prs)
    rect(s, 0, 0, 13.333, 7.5, AZUL_ESC)
    rect(s, 0, 5.0, 13.333, 2.5, AZUL)

    txt(s, "✅", 5.7, 0.6, 2.0, 1.5, size=64, color=VERDE, align=PP_ALIGN.CENTER)

    txt(s, "Sistema atualizado com sucesso!", 1.5, 2.0, 10.5, 0.85,
        size=32, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    txt(s, "As 7 novas funcionalidades já estão disponíveis para a Lilian e seus pacientes.",
        1.5, 2.85, 10.5, 0.55, size=16, color=RGBColor(0xC5, 0xDA, 0xF0), align=PP_ALIGN.CENTER)

    rect(s, 2.5, 3.55, 8.5, 0.05, VERDE)

    resumo = [
        "📅  Data de nascimento aprimorada",
        "🏷️  Novos status de paciente",
        "🔁  Agendamentos recorrentes",
        "💬  Lembretes por WhatsApp",
        "🎂  Mensagens de aniversário",
        "📝  Contratos com assinatura digital",
        "📋  Aba de Laudo no prontuário",
    ]

    for i, item in enumerate(resumo):
        col = i % 4
        row = i // 4
        x = 0.8 + col * 3.15
        y = 3.75 + row * 0.65
        txt(s, item, x, y, 3.0, 0.55, size=12, color=BRANCO)

    txt(s, "PsicoDoc v2.0 · Abril 2026", 0.5, 6.9, 12.5, 0.45,
        size=13, color=RGBColor(0x80, 0xA8, 0xCC), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════

def main():
    prs = nova_apresentacao()
    TOTAL = 9

    slide_capa(prs)
    slide_novidades(prs, 2, TOTAL)
    slide_visao_geral(prs, 3, TOTAL)
    slide_etapa1_sql(prs, 4, TOTAL)
    slide_etapa2_github(prs, 5, TOTAL)
    slide_etapa3_vercel(prs, 6, TOTAL)
    slide_testar(prs, 7, TOTAL)
    slide_whatsapp_config(prs, 8, TOTAL)
    slide_problemas(prs, 9, TOTAL)
    slide_final(prs, TOTAL, TOTAL)

    output = "PsicoDoc_Guia_Atualizacao_v2.pptx"
    prs.save(output)
    print(f"✅  Apresentação gerada: {output}")

if __name__ == "__main__":
    main()
